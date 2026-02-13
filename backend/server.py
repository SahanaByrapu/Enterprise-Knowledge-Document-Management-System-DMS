from fastapi import FastAPI, APIRouter, HTTPException, Depends, UploadFile, File, Query
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone, timedelta
import jwt
import bcrypt
import io
import json
import numpy as np
from emergentintegrations.llm.chat import LlmChat, UserMessage

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Emergent LLM Key
EMERGENT_LLM_KEY = os.environ.get('EMERGENT_LLM_KEY')

# JWT settings
JWT_SECRET = os.environ.get('JWT_SECRET', 'default-secret-key')
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_HOURS = 24

# Create the main app
app = FastAPI(title="Enterprise Knowledge DMS API")
api_router = APIRouter(prefix="/api")
security = HTTPBearer()

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# =========================
# MODELS
# =========================

class UserCreate(BaseModel):
    email: EmailStr
    password: str
    name: str

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class UserResponse(BaseModel):
    id: str
    email: str
    name: str
    role: str
    created_at: str

class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    user: UserResponse

class DocumentUpload(BaseModel):
    filename: str
    content_type: str
    size: int
    status: str = "pending"

class DocumentResponse(BaseModel):
    id: str
    filename: str
    content_type: str
    size: int
    status: str
    uploaded_by: str
    uploaded_at: str
    chunk_count: int = 0

class SearchQuery(BaseModel):
    query: str
    limit: int = 10

class SearchResult(BaseModel):
    document_id: str
    filename: str
    chunk_text: str
    score: float

class ChatMessage(BaseModel):
    message: str
    session_id: Optional[str] = None

class ChatResponse(BaseModel):
    response: str
    session_id: str
    sources: List[Dict[str, Any]]

class RoleUpdate(BaseModel):
    user_id: str
    role: str

class AuditLogResponse(BaseModel):
    id: str
    user_id: str
    user_email: str
    action: str
    resource: str
    details: str
    timestamp: str

class AnalyticsResponse(BaseModel):
    total_documents: int
    total_users: int
    total_queries: int
    access_trends: List[Dict[str, Any]]
    top_queries: List[Dict[str, Any]]

# =========================
# UTILITIES
# =========================

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str, email: str, role: str) -> str:
    payload = {
        "sub": user_id,
        "email": email,
        "role": role,
        "exp": datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    try:
        payload = jwt.decode(credentials.credentials, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user = await db.users.find_one({"id": payload["sub"]}, {"_id": 0})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

async def require_admin(user: dict = Depends(get_current_user)):
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return user

async def log_audit(user_id: str, user_email: str, action: str, resource: str, details: str):
    audit_log = {
        "id": str(uuid.uuid4()),
        "user_id": user_id,
        "user_email": user_email,
        "action": action,
        "resource": resource,
        "details": details,
        "timestamp": datetime.now(timezone.utc).isoformat()
    }
    await db.audit_logs.insert_one(audit_log)

def extract_text_from_file(file_content: bytes, content_type: str, filename: str) -> str:
    """Extract text from various document formats"""
    text = ""
    try:
        if content_type == "application/pdf" or filename.endswith('.pdf'):
            from PyPDF2 import PdfReader
            pdf = PdfReader(io.BytesIO(file_content))
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.endswith('.docx'):
            from docx import Document
            doc = Document(io.BytesIO(file_content))
            for para in doc.paragraphs:
                text += para.text + "\n"
        
        elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or filename.endswith('.xlsx'):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_content), data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"
        
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or filename.endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        elif content_type == "text/plain" or filename.endswith('.txt'):
            text = file_content.decode('utf-8', errors='ignore')
        
        else:
            text = file_content.decode('utf-8', errors='ignore')
    
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}")
        text = ""
    
    return text.strip()

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks"""
    if not text:
        return []
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    
    return chunks

def get_text_keywords(text: str) -> List[str]:
    """Extract keywords from text for simple TF-IDF style matching"""
    import re
    # Simple word tokenization and normalization
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    # Remove common stop words
    stop_words = {'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'her', 'was', 'one', 'our', 'out', 'has', 'have', 'been', 'from', 'they', 'were', 'been', 'have', 'many', 'some', 'them', 'this', 'that', 'with', 'will', 'would', 'there', 'their', 'what', 'which', 'when', 'where', 'while', 'into', 'more', 'other', 'could', 'should'}
    return [w for w in words if w not in stop_words]

def calculate_text_similarity(text1: str, text2: str) -> float:
    """Calculate similarity between two texts using keyword overlap"""
    words1 = set(get_text_keywords(text1))
    words2 = set(get_text_keywords(text2))
    
    if not words1 or not words2:
        return 0.0
    
    intersection = words1.intersection(words2)
    union = words1.union(words2)
    
    # Jaccard similarity
    return len(intersection) / len(union) if union else 0.0

def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    """Calculate cosine similarity between two vectors"""
    if not vec1 or not vec2:
        return 0.0
    a = np.array(vec1)
    b = np.array(vec2)
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

# =========================
# AUTH ROUTES
# =========================

@api_router.post("/auth/register", response_model=TokenResponse)
async def register(user_data: UserCreate):
    existing = await db.users.find_one({"email": user_data.email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    user_id = str(uuid.uuid4())
    user = {
        "id": user_id,
        "email": user_data.email,
        "password": hash_password(user_data.password),
        "name": user_data.name,
        "role": "user",
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.users.insert_one(user)
    
    token = create_token(user_id, user_data.email, "user")
    await log_audit(user_id, user_data.email, "REGISTER", "user", f"New user registered: {user_data.name}")
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user_id,
            email=user_data.email,
            name=user_data.name,
            role="user",
            created_at=user["created_at"]
        )
    )

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(credentials: UserLogin):
    user = await db.users.find_one({"email": credentials.email}, {"_id": 0})
    if not user or not verify_password(credentials.password, user["password"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    
    token = create_token(user["id"], user["email"], user["role"])
    await log_audit(user["id"], user["email"], "LOGIN", "auth", "User logged in")
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user["id"],
            email=user["email"],
            name=user["name"],
            role=user["role"],
            created_at=user["created_at"]
        )
    )

@api_router.get("/auth/me", response_model=UserResponse)
async def get_me(user: dict = Depends(get_current_user)):
    return UserResponse(
        id=user["id"],
        email=user["email"],
        name=user["name"],
        role=user["role"],
        created_at=user["created_at"]
    )

# =========================
# DOCUMENT ROUTES
# =========================

@api_router.post("/documents/upload", response_model=DocumentResponse)
async def upload_document(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    allowed_types = [
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "text/plain"
    ]
    
    # Check by extension if content_type is generic
    allowed_extensions = ['.pdf', '.docx', '.xlsx', '.pptx', '.txt']
    file_ext = '.' + file.filename.split('.')[-1].lower() if '.' in file.filename else ''
    
    if file.content_type not in allowed_types and file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail="File type not supported")
    
    content = await file.read()
    doc_id = str(uuid.uuid4())
    
    document = {
        "id": doc_id,
        "filename": file.filename,
        "content_type": file.content_type,
        "size": len(content),
        "status": "processing",
        "uploaded_by": user["id"],
        "uploaded_by_email": user["email"],
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "chunk_count": 0,
        "content": content  # Store raw content for processing
    }
    await db.documents.insert_one(document)
    
    # Process document asynchronously
    try:
        text = extract_text_from_file(content, file.content_type, file.filename)
        chunks = chunk_text(text)
        
        # Store chunks (text-based search, no embeddings needed)
        for i, chunk in enumerate(chunks):
            chunk_doc = {
                "id": str(uuid.uuid4()),
                "document_id": doc_id,
                "chunk_index": i,
                "text": chunk,
                "keywords": get_text_keywords(chunk),
                "created_at": datetime.now(timezone.utc).isoformat()
            }
            await db.document_chunks.insert_one(chunk_doc)
        
        # Update document status
        await db.documents.update_one(
            {"id": doc_id},
            {"$set": {"status": "indexed", "chunk_count": len(chunks)}}
        )
        
        await log_audit(user["id"], user["email"], "UPLOAD", "document", f"Uploaded: {file.filename}")
        
        return DocumentResponse(
            id=doc_id,
            filename=file.filename,
            content_type=file.content_type,
            size=len(content),
            status="indexed",
            uploaded_by=user["id"],
            uploaded_at=document["uploaded_at"],
            chunk_count=len(chunks)
        )
    
    except Exception as e:
        logger.error(f"Error processing document: {str(e)}")
        await db.documents.update_one(
            {"id": doc_id},
            {"$set": {"status": "failed"}}
        )
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")

@api_router.get("/documents", response_model=List[DocumentResponse])
async def get_documents(user: dict = Depends(get_current_user)):
    documents = await db.documents.find(
        {},
        {"_id": 0, "content": 0}
    ).sort("uploaded_at", -1).to_list(1000)
    
    return [
        DocumentResponse(
            id=doc["id"],
            filename=doc["filename"],
            content_type=doc["content_type"],
            size=doc["size"],
            status=doc["status"],
            uploaded_by=doc["uploaded_by"],
            uploaded_at=doc["uploaded_at"],
            chunk_count=doc.get("chunk_count", 0)
        )
        for doc in documents
    ]

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str, user: dict = Depends(get_current_user)):
    doc = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Only admin or owner can delete
    if user["role"] != "admin" and doc["uploaded_by"] != user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized to delete this document")
    
    await db.documents.delete_one({"id": doc_id})
    await db.document_chunks.delete_many({"document_id": doc_id})
    
    await log_audit(user["id"], user["email"], "DELETE", "document", f"Deleted: {doc['filename']}")
    
    return {"message": "Document deleted successfully"}

# =========================
# SEARCH ROUTES
# =========================

@api_router.post("/search", response_model=List[SearchResult])
async def semantic_search(query: SearchQuery, user: dict = Depends(get_current_user)):
    # Log the search query
    await db.search_queries.insert_one({
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "query": query.query,
        "timestamp": datetime.now(timezone.utc).isoformat()
    })
    
    # Get all chunks and calculate similarity using text matching
    chunks = await db.document_chunks.find({}, {"_id": 0}).to_list(10000)
    
    results = []
    for chunk in chunks:
        if chunk.get("text"):
            score = calculate_text_similarity(query.query, chunk["text"])
            if score > 0:  # Only include chunks with some match
                results.append({
                    "document_id": chunk["document_id"],
                    "chunk_text": chunk["text"],
                    "score": score
                })
    
    # Sort by score and limit
    results.sort(key=lambda x: x["score"], reverse=True)
    results = results[:query.limit]
    
    # Get document filenames
    final_results = []
    for result in results:
        doc = await db.documents.find_one({"id": result["document_id"]}, {"_id": 0, "filename": 1})
        if doc:
            final_results.append(SearchResult(
                document_id=result["document_id"],
                filename=doc["filename"],
                chunk_text=result["chunk_text"][:500],
                score=result["score"]
            ))
    
    await log_audit(user["id"], user["email"], "SEARCH", "query", f"Search: {query.query[:50]}")
    
    return final_results

# =========================
# RAG CHAT ROUTES
# =========================

@api_router.post("/chat", response_model=ChatResponse)
async def rag_chat(message: ChatMessage, user: dict = Depends(get_current_user)):
    session_id = message.session_id or str(uuid.uuid4())
    
    # Get relevant context using text-based search
    chunks = await db.document_chunks.find({}, {"_id": 0}).to_list(10000)
    
    # Find most relevant chunks using text similarity
    scored_chunks = []
    for chunk in chunks:
        if chunk.get("text"):
            score = calculate_text_similarity(message.message, chunk["text"])
            scored_chunks.append({
                "document_id": chunk["document_id"],
                "text": chunk["text"],
                "score": score
            })
    
    scored_chunks.sort(key=lambda x: x["score"], reverse=True)
    top_chunks = scored_chunks[:5]
    
    # Build context
    context = "\n\n".join([f"[Document excerpt]: {chunk['text']}" for chunk in top_chunks if chunk['score'] > 0])
    
    if not context:
        context = "No relevant documents found in the knowledge base."
    
    # Get chat history for session
    history = await db.chat_history.find(
        {"session_id": session_id},
        {"_id": 0}
    ).sort("timestamp", 1).to_list(20)
    
    system_message = f"""You are an intelligent knowledge assistant for an enterprise document management system. 
Your task is to answer questions based on the provided document context. 
If the answer is not found in the context, say so clearly.
Always cite which document the information comes from when possible.

Context from documents:
{context}"""
    
    # Use Emergent LLM integration
    try:
        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=f"rag-{session_id}",
            system_message=system_message
        ).with_model("openai", "gpt-4o")
        
        # Build conversation with history
        full_conversation = ""
        for h in history[-10:]:
            if h["role"] == "user":
                full_conversation += f"User: {h['content']}\n"
            else:
                full_conversation += f"Assistant: {h['content']}\n"
        
        full_conversation += f"User: {message.message}"
        
        user_message = UserMessage(text=full_conversation if full_conversation else message.message)
        assistant_message = await chat.send_message(user_message)
        
    except Exception as e:
        logger.error(f"LLM API error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error generating response: {str(e)}")
    
    # Save to history
    timestamp = datetime.now(timezone.utc).isoformat()
    await db.chat_history.insert_one({
        "id": str(uuid.uuid4()),
        "session_id": session_id,
        "user_id": user["id"],
        "role": "user",
        "content": message.message,
        "timestamp": timestamp
    })
    await db.chat_history.insert_one({
        "id": str(uuid.uuid4()),
        "session_id": session_id,
        "user_id": user["id"],
        "role": "assistant",
        "content": assistant_message,
        "timestamp": timestamp
    })
    
    # Log query
    await db.search_queries.insert_one({
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "query": message.message,
        "type": "chat",
        "timestamp": timestamp
    })
    
    # Get source documents
    sources = []
    for chunk in top_chunks[:3]:
        if chunk["score"] > 0:
            doc = await db.documents.find_one({"id": chunk["document_id"]}, {"_id": 0, "filename": 1})
            if doc:
                sources.append({
                    "document_id": chunk["document_id"],
                    "filename": doc["filename"],
                    "relevance": round(chunk["score"], 3)
                })
    
    return ChatResponse(
        response=assistant_message,
        session_id=session_id,
        sources=sources
    )

@api_router.get("/chat/history/{session_id}")
async def get_chat_history(session_id: str, user: dict = Depends(get_current_user)):
    history = await db.chat_history.find(
        {"session_id": session_id, "user_id": user["id"]},
        {"_id": 0}
    ).sort("timestamp", 1).to_list(100)
    return history

# =========================
# ANALYTICS ROUTES
# =========================

@api_router.get("/analytics", response_model=AnalyticsResponse)
async def get_analytics(user: dict = Depends(get_current_user)):
    # Basic counts
    total_documents = await db.documents.count_documents({})
    total_users = await db.users.count_documents({})
    total_queries = await db.search_queries.count_documents({})
    
    # Access trends (last 7 days)
    seven_days_ago = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
    
    # Aggregate queries by day
    pipeline = [
        {"$match": {"timestamp": {"$gte": seven_days_ago}}},
        {"$project": {
            "day": {"$substr": ["$timestamp", 0, 10]}
        }},
        {"$group": {"_id": "$day", "count": {"$sum": 1}}},
        {"$sort": {"_id": 1}}
    ]
    access_trends_raw = await db.search_queries.aggregate(pipeline).to_list(100)
    access_trends = [{"date": item["_id"], "queries": item["count"]} for item in access_trends_raw]
    
    # Fill in missing days
    all_days = []
    for i in range(7):
        day = (datetime.now(timezone.utc) - timedelta(days=6-i)).strftime("%Y-%m-%d")
        existing = next((t for t in access_trends if t["date"] == day), None)
        all_days.append({"date": day, "queries": existing["queries"] if existing else 0})
    
    # Top queries
    top_pipeline = [
        {"$group": {"_id": "$query", "count": {"$sum": 1}}},
        {"$sort": {"count": -1}},
        {"$limit": 10}
    ]
    top_queries_raw = await db.search_queries.aggregate(top_pipeline).to_list(10)
    top_queries = [{"query": item["_id"], "count": item["count"]} for item in top_queries_raw]
    
    return AnalyticsResponse(
        total_documents=total_documents,
        total_users=total_users,
        total_queries=total_queries,
        access_trends=all_days,
        top_queries=top_queries
    )

# =========================
# ADMIN ROUTES
# =========================

@api_router.get("/admin/users", response_model=List[UserResponse])
async def get_all_users(admin: dict = Depends(require_admin)):
    users = await db.users.find({}, {"_id": 0, "password": 0}).to_list(1000)
    return [
        UserResponse(
            id=u["id"],
            email=u["email"],
            name=u["name"],
            role=u["role"],
            created_at=u["created_at"]
        )
        for u in users
    ]

@api_router.put("/admin/users/role")
async def update_user_role(role_update: RoleUpdate, admin: dict = Depends(require_admin)):
    if role_update.role not in ["user", "admin", "viewer"]:
        raise HTTPException(status_code=400, detail="Invalid role")
    
    result = await db.users.update_one(
        {"id": role_update.user_id},
        {"$set": {"role": role_update.role}}
    )
    
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="User not found")
    
    await log_audit(admin["id"], admin["email"], "UPDATE_ROLE", "user", f"Changed user {role_update.user_id} role to {role_update.role}")
    
    return {"message": "Role updated successfully"}

@api_router.get("/admin/audit-logs", response_model=List[AuditLogResponse])
async def get_audit_logs(
    limit: int = Query(default=100, le=500),
    admin: dict = Depends(require_admin)
):
    logs = await db.audit_logs.find({}, {"_id": 0}).sort("timestamp", -1).to_list(limit)
    return [
        AuditLogResponse(
            id=log["id"],
            user_id=log["user_id"],
            user_email=log["user_email"],
            action=log["action"],
            resource=log["resource"],
            details=log["details"],
            timestamp=log["timestamp"]
        )
        for log in logs
    ]

@api_router.delete("/admin/users/{user_id}")
async def delete_user(user_id: str, admin: dict = Depends(require_admin)):
    if user_id == admin["id"]:
        raise HTTPException(status_code=400, detail="Cannot delete yourself")
    
    user = await db.users.find_one({"id": user_id}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    await db.users.delete_one({"id": user_id})
    await log_audit(admin["id"], admin["email"], "DELETE_USER", "user", f"Deleted user: {user['email']}")
    
    return {"message": "User deleted successfully"}

# =========================
# HEALTH CHECK
# =========================

@api_router.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now(timezone.utc).isoformat()}

# Include router
app.include_router(api_router)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
